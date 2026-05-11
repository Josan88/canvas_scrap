import os
from typing import Tuple

from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}

PAGE_SIZE = LETTER
MARGIN = 0.75 * inch


def is_office_file(filename: str) -> bool:
    """Return True if the filename has an Office extension we can convert."""
    return os.path.splitext(filename)[1].lower() in OFFICE_EXTENSIONS


# ---------------------------------------------------------------------------
# docx -> PDF
# ---------------------------------------------------------------------------


def _docx_to_pdf(input_path: str, output_pdf_path: str) -> None:
    from docx import Document
    from docx.table import Table as DocxTable

    doc = Document(input_path)
    styles = getSampleStyleSheet()
    story = []

    heading_sizes = {1: 20, 2: 16, 3: 14, 4: 12}

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            para = _docx_paragraph_from_element(element, doc, styles, heading_sizes)
            if para:
                story.append(para)
                story.append(Spacer(1, 4))

        elif tag == "tbl":
            table = _docx_table_from_element(element, doc, styles)
            if table:
                story.append(table)
                story.append(Spacer(1, 8))

    if not story:
        story.append(Paragraph("(Empty document)", styles["Normal"]))

    pdf = SimpleDocTemplate(
        output_pdf_path,
        pagesize=PAGE_SIZE,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN,
    )
    pdf.build(story)


def _docx_paragraph_from_element(element, doc, styles, heading_sizes):
    from docx.oxml.ns import qn

    style_name = None
    pPr = element.find(qn("w:pPr"))
    if pPr is not None:
        pStyle = pPr.find(qn("w:pStyle"))
        if pStyle is not None:
            style_name = pStyle.get(qn("w:val"), "")

    is_heading = style_name and style_name.startswith("Heading")
    level = 0
    if is_heading:
        try:
            level = int(style_name.replace("Heading", "").strip())
        except ValueError:
            level = 1

    runs_data = []
    for r in element.iter(qn("w:r")):
        text_el = r.find(qn("w:t"))
        if text_el is None or text_el.text is None:
            continue
        text = text_el.text
        bold = False
        italic = False
        rPr = r.find(qn("w:rPr"))
        if rPr is not None:
            if rPr.find(qn("w:b")) is not None:
                bold_el = rPr.find(qn("w:b"))
                val = bold_el.get(qn("w:val"), "true") if bold_el is not None else "true"
                bold = val.lower() in ("true", "1", "on")
            if rPr.find(qn("w:i")) is not None:
                italic_el = rPr.find(qn("w:i"))
                val = italic_el.get(qn("w:val"), "true") if italic_el is not None else "true"
                italic = val.lower() in ("true", "1", "on")
        runs_data.append((text, bold, italic))

    full_text = "".join(t for t, _, _ in runs_data)
    if not full_text.strip():
        return None

    # Build a single Paragraph with inline bold/italic markup
    safe_parts = []
    for text, bold, italic in runs_data:
        part = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if bold:
            part = f"<b>{part}</b>"
        if italic:
            part = f"<i>{part}</i>"
        safe_parts.append(part)
    combined = "".join(safe_parts)

    if is_heading:
        font_size = heading_sizes.get(level, 12)
        style = ParagraphStyle(
            f"Heading{level}",
            parent=styles["Normal"],
            fontSize=font_size,
            leading=font_size + 4,
            spaceAfter=6,
            textColor=colors.HexColor("#1a1a2e"),
        )
    else:
        style = styles["Normal"]

    return Paragraph(combined, style)


def _docx_table_from_element(element, doc, styles):
    from docx.oxml.ns import qn

    rows = element.findall(qn("w:tr"))
    if not rows:
        return None

    table_data = []
    for row in rows:
        cells = row.findall(qn("w:tc"))
        row_data = []
        for cell in cells:
            cell_text_parts = []
            for p in cell.findall(qn("w:p")):
                for r in p.iter(qn("w:r")):
                    t = r.find(qn("w:t"))
                    if t is not None and t.text:
                        cell_text_parts.append(t.text)
            row_data.append(" ".join(cell_text_parts) if cell_text_parts else "")
        table_data.append(row_data)

    if not table_data:
        return None

    col_count = max(len(r) for r in table_data)
    for row in table_data:
        while len(row) < col_count:
            row.append("")

    available_width = PAGE_SIZE[0] - 2 * MARGIN
    col_width = available_width / col_count if col_count else available_width

    table = Table(table_data, colWidths=[col_width] * col_count)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e0e0e0")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    return table


# ---------------------------------------------------------------------------
# pptx -> PDF
# ---------------------------------------------------------------------------


def _pptx_to_pdf(input_path: str, output_pdf_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(input_path)
    styles = getSampleStyleSheet()

    pdf = SimpleDocTemplate(
        output_pdf_path,
        pagesize=PAGE_SIZE,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN,
    )

    story = []
    slide_width = prs.slide_width or Emu(9144000)
    slide_height = prs.slide_height or Emu(6858000)

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx > 0:
            story.append(Spacer(1, 12))

        slide_label = ParagraphStyle(
            "SlideLabel",
            parent=styles["Normal"],
            fontSize=8,
            textColor=colors.grey,
        )
        story.append(Paragraph(f"--- Slide {slide_idx + 1} ---", slide_label))
        story.append(Spacer(1, 4))

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                bold_runs = any(r.font.bold for r in para.runs if r.font.bold)
                if bold_runs:
                    safe = f"<b>{safe}</b>"

                font_size = 11
                for run in para.runs:
                    if run.font.size:
                        font_size = max(8, min(run.font.size / 12700, 24))
                        break

                style = ParagraphStyle(
                    "SlideText",
                    parent=styles["Normal"],
                    fontSize=font_size,
                    leading=font_size + 3,
                )
                story.append(Paragraph(safe, style))
                story.append(Spacer(1, 2))

            if shape.has_table:
                tbl = shape.table
                table_data = []
                for row in tbl.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)

                if table_data:
                    col_count = max(len(r) for r in table_data)
                    available = PAGE_SIZE[0] - 2 * MARGIN
                    cw = available / col_count if col_count else available
                    t = Table(table_data, colWidths=[cw] * col_count)
                    t.setStyle(
                        TableStyle(
                            [
                                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                                ("FONTSIZE", (0, 0), (-1, -1), 8),
                                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
                                ("TOPPADDING", (0, 0), (-1, -1), 2),
                                ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
                            ]
                        )
                    )
                    story.append(t)
                    story.append(Spacer(1, 4))

    if not story:
        story.append(Paragraph("(Empty presentation)", styles["Normal"]))

    pdf.build(story)


# ---------------------------------------------------------------------------
# xlsx -> PDF
# ---------------------------------------------------------------------------


def _xlsx_to_pdf(input_path: str, output_pdf_path: str) -> None:
    from openpyxl import load_workbook

    wb = load_workbook(input_path, read_only=True, data_only=True)
    styles = getSampleStyleSheet()

    pdf = SimpleDocTemplate(
        output_pdf_path,
        pagesize=PAGE_SIZE,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN,
    )

    story = []

    for sheet_idx, ws in enumerate(wb.worksheets):
        if sheet_idx > 0:
            story.append(Spacer(1, 16))

        title_style = ParagraphStyle(
            "SheetTitle",
            parent=styles["Normal"],
            fontSize=14,
            leading=18,
            spaceAfter=6,
            textColor=colors.HexColor("#1a1a2e"),
        )
        story.append(Paragraph(ws.title, title_style))
        story.append(Spacer(1, 4))

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            story.append(Paragraph("(Empty sheet)", styles["Normal"]))
            continue

        # Trim trailing empty rows
        while rows and all(c is None or str(c).strip() == "" for c in rows[-1]):
            rows.pop()
        if not rows:
            story.append(Paragraph("(Empty sheet)", styles["Normal"]))
            continue

        col_count = max(len(r) for r in rows)
        for row in rows:
            while len(row) < col_count:
                row.append(None)

        def _cell_str(val):
            if val is None:
                return ""
            return str(val).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        table_data = [[_cell_str(c) for c in row] for row in rows]

        available = PAGE_SIZE[0] - 2 * MARGIN
        col_width = available / col_count if col_count else available

        t = Table(table_data, colWidths=[col_width] * col_count)
        style_cmds = [
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3),
            ("TOPPADDING", (0, 0), (-1, -1), 2),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
        ]
        if len(rows) > 0:
            style_cmds.append(("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e0e0e0")))
        t.setStyle(TableStyle(style_cmds))
        story.append(t)

    if not story:
        story.append(Paragraph("(Empty workbook)", styles["Normal"]))

    wb.close()
    pdf.build(story)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def convert_office_to_pdf(
    input_path: str, output_dir: str
) -> Tuple[bool, str, str]:
    """Convert an Office file to PDF.

    Returns:
        (success, pdf_path, error_message)
    """
    ext = os.path.splitext(input_path)[1].lower()
    base = os.path.splitext(os.path.basename(input_path))[0]
    pdf_path = os.path.join(output_dir, f"{base}.pdf")

    try:
        if ext == ".docx":
            _docx_to_pdf(input_path, pdf_path)
        elif ext == ".pptx":
            _pptx_to_pdf(input_path, pdf_path)
        elif ext == ".xlsx":
            _xlsx_to_pdf(input_path, pdf_path)
        else:
            return False, "", f"Unsupported office extension: {ext}"

        if not os.path.exists(pdf_path):
            return False, "", "PDF file was not created"
        return True, pdf_path, ""

    except Exception as e:
        return False, "", f"Office-to-PDF conversion failed: {e}"
